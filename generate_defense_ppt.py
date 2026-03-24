#!/usr/bin/env python3
"""
生成硕士论文答辩PPT
论文题目：基于GA-HPO FE-IDDQN的云工作流调度优化研究
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
import os

# ==================== 配色方案 ====================
PRIMARY = RGBColor(0x1A, 0x3C, 0x6E)      # 深蓝主色
ACCENT = RGBColor(0x2E, 0x86, 0xC1)       # 亮蓝强调
ACCENT2 = RGBColor(0x27, 0xAE, 0x60)      # 绿色
ACCENT3 = RGBColor(0xE7, 0x4C, 0x3C)      # 红色
DARK = RGBColor(0x2C, 0x3E, 0x50)         # 深灰文字
LIGHT = RGBColor(0xEC, 0xF0, 0xF1)        # 浅灰背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
GOLD = RGBColor(0xF3, 0x9C, 0x12)         # 金色
LIGHT_BLUE_BG = RGBColor(0xD6, 0xEA, 0xF8)
VERY_LIGHT_GREY = RGBColor(0xF8, 0xF9, 0xF9)


def add_slide_number(slide, prs, num):
    """在右下角添加页码"""
    left = prs.slide_width - Inches(0.8)
    top = prs.slide_height - Inches(0.5)
    txBox = slide.shapes.add_textbox(left, top, Inches(0.6), Inches(0.3))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(num)
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.RIGHT


def add_bottom_bar(slide, prs):
    """底部装饰条"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), prs.slide_height - Inches(0.08),
        prs.slide_width, Inches(0.08)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = PRIMARY
    bar.line.fill.background()


def add_top_bar(slide, prs, color=PRIMARY):
    """顶部装饰条"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(0.06)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def set_slide_bg(slide, color=WHITE):
    """设置幻灯片背景色"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_title_shape(slide, text, left, top, width, height, font_size=28, color=PRIMARY, bold=True, alignment=PP_ALIGN.LEFT):
    """添加标题文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


def add_body_text(slide, text, left, top, width, height, font_size=16, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, line_spacing=1.5):
    """添加正文文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_bullet_list(slide, items, left, top, width, height, font_size=15, color=DARK, bullet_color=ACCENT, spacing=1.3):
    """添加列表项"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # 使用 bullet 符号
        if isinstance(item, tuple):
            # (title, detail)
            run1 = p.add_run()
            run1.text = "● " + item[0]
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = bullet_color
            run1.font.bold = True
            run2 = p.add_run()
            run2.text = " — " + item[1]
            run2.font.size = Pt(font_size - 1)
            run2.font.color.rgb = color
        else:
            p.text = "● " + item
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
        p.space_after = Pt(font_size * (spacing - 1) + 2)
    return txBox


def add_card(slide, title, content, left, top, width, height, title_color=ACCENT, bg_color=VERY_LIGHT_GREY):
    """添加卡片样式的内容块"""
    # 背景矩形
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = bg_color
    rect.line.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
    rect.line.width = Pt(0.5)

    # 顶部装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = title_color
    line.line.fill.background()

    # 标题
    add_title_shape(slide, title, left + Inches(0.2), top + Inches(0.18), width - Inches(0.4), Inches(0.4),
                    font_size=14, color=title_color, bold=True)
    # 内容
    add_body_text(slide, content, left + Inches(0.2), top + Inches(0.55), width - Inches(0.4), height - Inches(0.7),
                  font_size=12, color=DARK, line_spacing=1.3)


def add_table(slide, headers, rows, left, top, width, row_height=Inches(0.4)):
    """添加表格"""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    col_width = width // num_cols

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, row_height * num_rows)
    table = table_shape.table

    # 设置列宽
    for i in range(num_cols):
        table.columns[i].width = col_width

    # 表头
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 数据行
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = VERY_LIGHT_GREY
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = DARK
            p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_num = 0

    # ==================== Slide 1: 封面 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    set_slide_bg(slide, WHITE)

    # 顶部深蓝色背景块
    top_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, Inches(3.2))
    top_rect.fill.solid()
    top_rect.fill.fore_color.rgb = PRIMARY
    top_rect.line.fill.background()

    # 学校信息
    add_title_shape(slide, "硕士学位论文答辩", Inches(0.5), Inches(0.4), Inches(12), Inches(0.5),
                    font_size=18, color=RGBColor(0xBB, 0xCC, 0xDD), bold=False, alignment=PP_ALIGN.CENTER)

    # 论文标题
    add_title_shape(slide, "基于GA-HPO FE-IDDQN的\n云工作流调度优化研究", Inches(0.5), Inches(1.0), Inches(12), Inches(1.8),
                    font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 装饰线
    deco_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.5), Inches(2.3), Inches(0.04))
    deco_line.fill.solid()
    deco_line.fill.fore_color.rgb = ACCENT
    deco_line.line.fill.background()

    # 答辩人信息
    info_text = "答辩人：XXX\n指导教师：XXX 教授\n专    业：计算机技术\n日    期：2026年X月"
    add_body_text(slide, info_text, Inches(3.5), Inches(3.8), Inches(6), Inches(2.5),
                  font_size=20, color=DARK, alignment=PP_ALIGN.CENTER, line_spacing=1.8)

    add_bottom_bar(slide, prs)

    # ==================== Slide 2: 目录 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "汇 报 提 纲", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8),
                    font_size=32, color=PRIMARY, alignment=PP_ALIGN.CENTER)

    chapters = [
        ("01", "研究背景与意义", "云工作流调度挑战、现有方法不足"),
        ("02", "相关研究工作", "DRL调度研究现状、特征工程的必要性"),
        ("03", "问题建模与特征工程", "DAG工作流模型、多维特征体系、RL环境设计"),
        ("04", "GA-HPO FE-IDDQN算法设计", "双流网络、训练机制、超参数优化"),
        ("05", "实验验证与性能评估", "性能对比、消融实验、案例研究"),
        ("06", "总结与展望", "主要贡献、不足与未来方向"),
    ]

    for i, (num, title, desc) in enumerate(chapters):
        row = i // 2
        col = i % 2
        x = Inches(1.2 + col * 5.5)
        y = Inches(1.6 + row * 1.8)

        # 编号圆形
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y + Inches(0.05), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ACCENT if i < 5 else ACCENT2
        circle.line.fill.background()
        tf = circle.text_frame
        tf.paragraphs[0].text = num
        tf.paragraphs[0].font.size = Pt(18)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        add_title_shape(slide, title, x + Inches(0.8), y, Inches(4), Inches(0.4),
                        font_size=18, color=DARK, bold=True)
        add_body_text(slide, desc, x + Inches(0.8), y + Inches(0.4), Inches(4.5), Inches(0.35),
                      font_size=12, color=RGBColor(0x77, 0x77, 0x77), line_spacing=1.2)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 3: 研究背景 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "一、研究背景与意义", Inches(0.5), Inches(0.25), Inches(8), Inches(0.6),
                    font_size=28, color=PRIMARY)

    # 左侧：背景
    add_card(slide, "行业背景", "● 云计算环境下工作流调度是核心技术挑战\n● 金融行业DAG工作流复杂度高（最多555个任务）\n● DolphinScheduler广泛应用于企业级调度\n● 当前调度以启发式为主，缺乏智能优化",
             Inches(0.5), Inches(1.1), Inches(5.8), Inches(2.6), title_color=ACCENT)

    # 右侧：问题
    add_card(slide, "核心问题", "● 调度空间大：任务-资源分配组合呈指数增长\n● 资源异构性：Worker节点性能差异显著\n● 依赖约束复杂：DAG拓扑约束限制并行度\n● 多目标冲突：Makespan vs 资源利用率 vs 负载均衡",
             Inches(6.8), Inches(1.1), Inches(5.8), Inches(2.6), title_color=ACCENT3)

    # 底部：研究意义
    add_card(slide, "研究意义", "● 理论层面：探索DRL在DAG工作流调度中的有效性，提出特征增强与超参数自动优化的技术路线\n● 实践层面：基于金融行业DolphinScheduler生产数据验证，优化监管报送等关键业务流程",
             Inches(0.5), Inches(4.0), Inches(12.1), Inches(2.0), title_color=ACCENT2)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 4: 现有方法不足 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "一、现有方法的不足与本文定位", Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
                    font_size=28, color=PRIMARY)

    # 三列对比
    add_card(slide, "传统启发式", "✗ FIFO/SJF/EFT等规则固定\n✗ 无法适应动态环境\n✗ 难以全局优化\n✗ 缺乏对DAG结构的利用\n✓ 计算速度快（<0.1ms）",
             Inches(0.5), Inches(1.2), Inches(3.8), Inches(2.8), title_color=RGBColor(0x95, 0xA5, 0xA6))

    add_card(slide, "元启发式 (GA/PSO/ACO)", "✗ 计算开销大（秒级/次）\n✗ 不适用于在线实时调度\n✗ 超参数敏感\n✓ 全局搜索能力\n✓ 可处理组合优化",
             Inches(4.6), Inches(1.2), Inches(3.8), Inches(2.8), title_color=RGBColor(0x95, 0xA5, 0xA6))

    add_card(slide, "现有DRL方法", "✗ 特征工程粗糙，状态表示不充分\n✗ 超参数高度依赖人工经验\n✗ 冷启动问题突出\n✗ 缺乏对DAG拓扑的显式建模\n✓ 端到端学习，在线推理快",
             Inches(8.7), Inches(1.2), Inches(3.8), Inches(2.8), title_color=RGBColor(0x95, 0xA5, 0xA6))

    # 本文方案
    add_card(slide, "▶ 本文方案：GA-HPO FE-IDDQN",
             "特征增强（16+7维领域特征）+ 双流Dueling网络 + 交叉注意力融合 + Graph Transformer\n"
             "SJF专家预填充 + PER + 动作掩码 → 训练效率↑  收敛速度↑\n"
             "GA网络架构搜索 + Optuna训练超参搜索 → 自动化调优，消除人工经验依赖\n"
             "推理延迟~5ms，满足在线实时调度需求",
             Inches(0.5), Inches(4.3), Inches(12.1), Inches(2.5), title_color=ACCENT)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 5: 问题建模 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "二、问题建模与特征工程", Inches(0.5), Inches(0.25), Inches(8), Inches(0.6),
                    font_size=28, color=PRIMARY)

    # DAG工作流模型
    add_card(slide, "DAG工作流模型", "工作流 W = (T, E, R)：\n"
             "● T = {t₁, t₂, ..., tₙ}  任务集合\n"
             "● E ⊆ T×T  任务间依赖关系（有向无环图边集）\n"
             "● R = {r₁, r₂, ..., rₘ}  异构资源集合（5个Worker节点）\n"
             "目标：min Makespan = max end(tᵢ) - min start(tᵢ)",
             Inches(0.5), Inches(1.1), Inches(5.8), Inches(3.0), title_color=ACCENT)

    # 多目标优化
    add_card(slide, "多目标优化函数", "min F = w₁·Makespan + w₂·(1-Util) + w₃·LoadImbalance\n\n"
             "权重配置：\n"
             "● w₁ = 0.5  （Makespan最小化）\n"
             "● w₂ = 0.2  （资源利用率最大化）\n"
             "● w₃ = 0.3  （负载均衡最大化）\n\n"
             "约束：所有DAG依赖关系必须满足",
             Inches(6.8), Inches(1.1), Inches(5.8), Inches(3.0), title_color=ACCENT2)

    # 特征工程
    add_card(slide, "多维度特征工程体系",
             "任务级特征（16维）：类型编码(1) + 资源需求(2) + 执行统计(4) + DAG结构(4) + 时间(3) + 上下文(2)\n"
             "资源级特征（7维）：CPU利用率 + 内存利用率 + 任务数 + 平均执行时间 + 可用时间 + 队列长度 + 累计任务数\n"
             "DAG拓扑信息：邻接矩阵 A ∈ {0,1}ⁿˣⁿ，支持Graph Transformer的注意力掩码",
             Inches(0.5), Inches(4.4), Inches(12.1), Inches(2.3), title_color=GOLD)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 6: RL环境设计 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "二、强化学习环境设计", Inches(0.5), Inches(0.25), Inches(8), Inches(0.6),
                    font_size=28, color=PRIMARY)

    add_card(slide, "状态空间 S（双流表示）",
             "S = (X_task, X_resource, A_dag)\n"
             "● X_task ∈ ℝⁿˣ¹⁶  当前就绪任务的特征矩阵\n"
             "● X_resource ∈ ℝ⁵ˣ⁷  资源节点状态矩阵\n"
             "● A_dag ∈ {0,1}ⁿˣⁿ  DAG邻接矩阵\n"
             "特征经过Min-Max归一化至[0,1]",
             Inches(0.5), Inches(1.1), Inches(3.8), Inches(3.0), title_color=ACCENT)

    add_card(slide, "动作空间 A",
             "A = {1, 2, 3, 4, 5}\n"
             "选择一个Worker节点执行当前任务\n\n"
             "动作掩码机制：\n"
             "● 屏蔽资源不足的节点\n"
             "● 屏蔽故障节点\n"
             "● 保证调度可行性",
             Inches(4.6), Inches(1.1), Inches(3.8), Inches(3.0), title_color=ACCENT2)

    add_card(slide, "奖励函数 R（多目标）",
             "R = -[w₁·MakespanCost + w₂·(1-Util) + w₃·LoadImbalance]\n\n"
             "● MakespanCost：归一化任务完成时间\n"
             "● Util：当前资源利用率均值\n"
             "● LoadImbalance：节点负载标准差\n"
             "即时奖励，每步计算",
             Inches(8.7), Inches(1.1), Inches(3.8), Inches(3.0), title_color=ACCENT3)

    # 仿真环境
    add_card(slide, "历史数据重放仿真器（HistoricalReplaySimulator）",
             "● 从DolphinScheduler数据库加载真实工作流实例数据\n"
             "● 按DAG拓扑排序逐任务调度，维护5个Worker节点状态\n"
             "● 任务执行时间基于历史实际执行时间，保证仿真真实性\n"
             "● 支持CONDITIONS和DEPENDENT等特殊任务类型的条件依赖处理",
             Inches(0.5), Inches(4.4), Inches(12.1), Inches(2.3), title_color=GOLD)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 7: 算法总体框架 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "三、GA-HPO FE-IDDQN 算法总体框架", Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
                    font_size=28, color=PRIMARY)

    # 五大模块
    modules = [
        ("特征工程模块", "异构数据→结构化\n状态表示（16+7维）", ACCENT),
        ("双流网络模块", "任务流+资源流\n交叉注意力融合→Q值", RGBColor(0x8E, 0x44, 0xAD)),
        ("训练机制模块", "PER + SJF预填充\nε-贪心 + 动作掩码", ACCENT2),
        ("超参数优化模块", "GA架构搜索\nOptuna训练参数搜索", GOLD),
        ("奖励计算模块", "多目标奖励函数\nMakespan+Util+LBI", ACCENT3),
    ]

    for i, (title, desc, color) in enumerate(modules):
        x = Inches(0.5 + i * 2.5)
        y = Inches(1.2)

        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(2.2), Inches(2.0))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()

        add_title_shape(slide, title, x + Inches(0.1), y + Inches(0.2), Inches(2.0), Inches(0.5),
                        font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        add_body_text(slide, desc, x + Inches(0.1), y + Inches(0.8), Inches(2.0), Inches(1.0),
                      font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER, line_spacing=1.3)

    # 工作流程
    add_card(slide, "离线训练阶段",
             "数据划分（6:2:2分层抽样）→ SJF预填充经验缓冲 → 与仿真环境交互学习\n"
             "→ PER优先采样 + 双重Q网络更新 → 验证集早停 → 保存最优模型",
             Inches(0.5), Inches(3.6), Inches(5.8), Inches(1.8), title_color=ACCENT)

    add_card(slide, "在线推理阶段",
             "接收实时任务/资源状态 → 特征提取（16+7维）→ 双流网络前向传播\n"
             "→ 动作掩码过滤 → 贪心选择最优Worker → 输出调度决策（~5ms延迟）",
             Inches(6.8), Inches(3.6), Inches(5.8), Inches(1.8), title_color=ACCENT2)

    # GA-HPO两阶段
    add_card(slide, "GA-HPO 两阶段超参数自动优化",
             "第一阶段 GA架构搜索（种群12×10代）：搜索hidden_dim, fusion_dim, num_heads, transformer_layers, dropout, use_gnn\n"
             "第二阶段 Optuna超参搜索（TPE+50次试验）：搜索lr, gamma, tau, ε_decay, batch_size, buffer_size, PER α/β, gradient_clip",
             Inches(0.5), Inches(5.6), Inches(12.1), Inches(1.3), title_color=GOLD)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 8: 双流网络架构 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "三、双流Dueling网络架构", Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
                    font_size=28, color=PRIMARY)

    # 任务流
    add_card(slide, "任务流（Task Stream）",
             "输入：X_task ∈ ℝⁿˣ¹⁶\n"
             "MLP编码：16→512→256→128 (ReLU+Dropout)\n"
             "多头自注意力：4 heads, 捕捉任务间关联\n"
             "可选Graph Transformer：利用DAG邻接矩阵掩码\n"
             "输出：H_task ∈ ℝⁿˣ¹²⁸",
             Inches(0.5), Inches(1.1), Inches(5.2), Inches(2.6), title_color=ACCENT)

    # 资源流
    add_card(slide, "资源流（Resource Stream）",
             "输入：X_resource ∈ ℝ⁵ˣ⁷\n"
             "MLP编码：7→512→256→128 (ReLU+Dropout)\n"
             "多头自注意力：4 heads, 捕捉资源间竞争\n"
             "输出：H_resource ∈ ℝ⁵ˣ¹²⁸",
             Inches(6.3), Inches(1.1), Inches(5.2), Inches(2.6), title_color=ACCENT2)

    # 融合与输出
    add_card(slide, "交叉注意力融合 + Dueling输出",
             "交叉注意力：Q=H_task, K=V=H_resource → 任务表示动态融入资源信息 → F ∈ ℝⁿˣ²⁵⁶\n"
             "融合层：任务向量⊕资源均值→全局状态表征 → 256维融合特征\n"
             "Dueling分解：Q(s,a) = V(s) + A(s,a) - mean(A)  分离状态价值与动作优势\n"
             "动作掩码：将不可行动作的Q值设为-∞ → 保证100%依赖约束满足",
             Inches(0.5), Inches(4.0), Inches(12.1), Inches(2.8), title_color=RGBColor(0x8E, 0x44, 0xAD))

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 9: 训练机制 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "三、核心训练机制", Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
                    font_size=28, color=PRIMARY)

    add_card(slide, "① 双重Q网络（Double DQN）",
             "● 主网络选择动作：a* = argmax Q(s',a;θ)\n"
             "● 目标网络评估价值：y = r + γ·Q(s',a*;θ⁻)\n"
             "● 解耦选择与评估 → 缓解Q值高估问题\n"
             "● 目标网络Polyak软更新：θ⁻ ← τ·θ + (1-τ)·θ⁻, τ=0.005",
             Inches(0.5), Inches(1.1), Inches(3.8), Inches(2.8), title_color=ACCENT)

    add_card(slide, "② 优先经验回放 (PER)",
             "● 基于SumTree的O(log N)高效采样\n"
             "● 采样概率 P(i) ∝ |δᵢ|^α，α=0.6\n"
             "● 重要性采样权重 wᵢ = (N·P(i))^(-β)\n"
             "● β从0.4线性增长到1.0\n"
             "● 优先学习TD误差大的关键调度经验",
             Inches(4.6), Inches(1.1), Inches(3.8), Inches(2.8), title_color=ACCENT2)

    add_card(slide, "③ SJF专家经验预填充",
             "● 冷启动解决方案\n"
             "● 用SJF启发式运行工作流生成经验\n"
             "● 高优先级填充至经验缓冲区\n"
             "● 收敛速度提高46%\n  (280→150 episodes)",
             Inches(8.7), Inches(1.1), Inches(3.8), Inches(2.8), title_color=ACCENT3)

    # 底部：探索策略
    add_card(slide, "④ SJF偏置ε-贪心探索 + 动作掩码",
             "探索阶段（概率ε）：50%执行SJF启发式动作 + 50%随机探索 → 提高有效经验质量\n"
             "利用阶段（概率1-ε）：贪心选择Q值最大的动作 → ε从1.0按0.998衰减至0.05\n"
             "动作掩码：始终保证选择的动作满足资源约束和依赖约束 → DSR=100%",
             Inches(0.5), Inches(4.2), Inches(12.1), Inches(2.3), title_color=GOLD)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 10: 实验设置 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "四、实验设置", Inches(0.5), Inches(0.25), Inches(8), Inches(0.6),
                    font_size=28, color=PRIMARY)

    # 数据集
    add_card(slide, "实验数据集（DolphinScheduler生产数据）",
             "● 来源：某大型保险公司WhaleScheduler生产环境\n"
             "● 工作流实例：~2,400个（成功执行，state=7）\n"
             "● 任务实例：~50,000个\n"
             "● 依赖关系：~2,000条\n"
             "● Worker节点：5个\n"
             "● 数据划分：训练60% / 验证20% / 测试20%（分层抽样）",
             Inches(0.5), Inches(1.1), Inches(5.8), Inches(3.0), title_color=ACCENT)

    # 对比算法
    add_card(slide, "对比算法（3类9种）",
             "传统启发式：\n  FIFO / RoundRobin / SJF / EFT\n\n"
             "元启发式：\n  GA(种群20,代数50) / PSO(粒子10,迭代30) / ACO(蚂蚁10,迭代20)\n\n"
             "深度强化学习：\n  DQN / DDQN（隐层[128,64], lr=1e-3）",
             Inches(6.8), Inches(1.1), Inches(5.8), Inches(3.0), title_color=ACCENT2)

    # 评价指标
    add_card(slide, "评价指标",
             "① Makespan（完工时间）— 首要优化目标    ② 资源利用率 — 资源有效使用程度\n"
             "③ 负载均衡指数(LBI) — 工作量分配均匀度    ④ 依赖约束满足率 — 调度正确性\n"
             "⑤ 推理延迟 — 在线可用性    公平保证：所有算法在相同测试集上评估，10次独立运行取均值±标准差",
             Inches(0.5), Inches(4.4), Inches(12.1), Inches(2.3), title_color=GOLD)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 11: 性能对比 - Makespan ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "四、性能对比 — 不同规模工作流Makespan", Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
                    font_size=28, color=PRIMARY)

    # 表格
    headers = ["算法", "小型(1-10)", "中型(11-30)", "大型(31+)", "总体平均"]
    rows = [
        ["FIFO", "43,547", "16,730", "21,093", "23,416"],
        ["RoundRobin", "36,272", "10,862", "5,121", "15,053"],
        ["SJF", "42,229", "10,900", "4,230", "14,668"],
        ["EFT", "42,229", "10,642", "4,230", "14,544"],
        ["GA", "38,450", "10,235", "4,087", "13,842"],
        ["DQN", "37,856", "10,524", "4,345", "14,235"],
        ["DDQN", "37,245", "10,287", "4,198", "13,986"],
        ["FE-IDDQN", "35,780", "9,876", "3,892", "13,182"],
    ]
    add_table(slide, headers, rows, Inches(0.5), Inches(1.1), Inches(7.5), Inches(0.38))

    # 图表 - 总体对比柱状图
    chart_data = CategoryChartData()
    chart_data.categories = ['FIFO', 'RR', 'SJF', 'EFT', 'GA', 'DQN', 'DDQN', 'FE-IDDQN']
    chart_data.add_series('总体Makespan(秒)', (23416, 15053, 14668, 14544, 13842, 14235, 13986, 13182))

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(8.5), Inches(1.1), Inches(4.5), Inches(3.0),
        chart_data
    ).chart
    chart.has_legend = False
    chart.style = 2
    plot = chart.plots[0]
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ACCENT

    # 关键结论
    add_card(slide, "关键发现",
             "● FE-IDDQN总体Makespan 13,182秒：vs FIFO ↓43.7%，vs EFT ↓9.4%，vs DDQN ↓5.7%\n"
             "● 大型工作流优势最显著：FE-IDDQN vs FIFO ↓81.5%，通过端到端RL学习跨步最优策略\n"
             "● 小型工作流差异较小：调度空间有限，启发式已可产生较优方案",
             Inches(0.5), Inches(4.8), Inches(12.1), Inches(2.0), title_color=ACCENT2)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 12: 资源利用率与负载均衡 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "四、资源利用率与负载均衡对比", Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
                    font_size=28, color=PRIMARY)

    # 资源利用率柱状图
    chart_data1 = CategoryChartData()
    chart_data1.categories = ['FIFO', 'RR', 'SJF', 'EFT', 'GA', 'DQN', 'DDQN', 'FE-IDDQN']
    chart_data1.add_series('资源利用率(%)', (20.0, 45.2, 51.1, 53.1, 54.8, 48.7, 52.3, 58.6))

    chart1 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.1), Inches(5.8), Inches(3.2),
        chart_data1
    ).chart
    chart1.has_legend = False
    chart1.chart_title.has_text_frame = True
    chart1.chart_title.text_frame.paragraphs[0].text = "平均资源利用率 (%)"
    chart1.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart1.chart_title.text_frame.paragraphs[0].font.color.rgb = DARK
    plot1 = chart1.plots[0]
    series1 = plot1.series[0]
    series1.format.fill.solid()
    series1.format.fill.fore_color.rgb = ACCENT2

    # 负载均衡柱状图
    chart_data2 = CategoryChartData()
    chart_data2.categories = ['FIFO', 'RR', 'SJF', 'EFT', 'GA', 'DQN', 'DDQN', 'FE-IDDQN']
    chart_data2.add_series('负载均衡指数', (0.21, 0.85, 0.78, 0.76, 0.79, 0.72, 0.75, 0.87))

    chart2 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(6.8), Inches(1.1), Inches(5.8), Inches(3.2),
        chart_data2
    ).chart
    chart2.has_legend = False
    chart2.chart_title.has_text_frame = True
    chart2.chart_title.text_frame.paragraphs[0].text = "负载均衡指数 (LBI)"
    chart2.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart2.chart_title.text_frame.paragraphs[0].font.color.rgb = DARK
    plot2 = chart2.plots[0]
    series2 = plot2.series[0]
    series2.format.fill.solid()
    series2.format.fill.fore_color.rgb = GOLD

    # 分析
    add_card(slide, "分析",
             "● FE-IDDQN资源利用率58.6%（最高），vs FIFO的20% → 多目标奖励函数引导资源有效使用\n"
             "● FE-IDDQN负载均衡0.87（最高），甚至超过天然均匀分配的RoundRobin(0.85) → 奖励函数中LBI权重0.3的效果\n"
             "● FE-IDDQN是唯一在利用率和均衡两个指标上均排名前列的算法 → 多目标协同优化成功",
             Inches(0.5), Inches(4.6), Inches(12.1), Inches(2.2), title_color=ACCENT)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 13: GA-HPO效果 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "四、GA-HPO超参数优化效果", Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
                    font_size=28, color=PRIMARY)

    # HPO对比表
    headers_hpo = ["配置", "Makespan(秒)", "利用率", "LBI", "搜索耗时"]
    rows_hpo = [
        ["Manual（手动调参）", "14,856", "51.2%", "0.76", "—"],
        ["GA-Only（仅架构搜索）", "13,978", "54.5%", "0.80", "~4h"],
        ["Optuna-Only（仅超参搜索）", "14,124", "53.8%", "0.79", "~3h"],
        ["GA-HPO（两阶段优化）", "13,182", "58.6%", "0.87", "~7h"],
    ]
    add_table(slide, headers_hpo, rows_hpo, Inches(0.5), Inches(1.1), Inches(7.0), Inches(0.42))

    # 柱状图对比
    chart_data3 = CategoryChartData()
    chart_data3.categories = ['Manual', 'GA-Only', 'Optuna-Only', 'GA-HPO']
    chart_data3.add_series('Makespan(秒)', (14856, 13978, 14124, 13182))

    chart3 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(8.0), Inches(1.1), Inches(4.8), Inches(2.5),
        chart_data3
    ).chart
    chart3.has_legend = False
    chart3.chart_title.has_text_frame = True
    chart3.chart_title.text_frame.paragraphs[0].text = "Makespan优化对比"
    chart3.chart_title.text_frame.paragraphs[0].font.size = Pt(11)
    plot3 = chart3.plots[0]
    series3 = plot3.series[0]
    series3.format.fill.solid()
    series3.format.fill.fore_color.rgb = GOLD

    # 分析
    add_card(slide, "关键结论",
             "● GA-HPO两阶段联合优化使Makespan降低11.3%（14,856→13,182秒）\n"
             "● GA架构搜索贡献5.9%，Optuna超参搜索贡献4.9%，两者存在协同效应\n"
             "● 搜索开销~7h为一次性离线成本，最优超参数可在后续训练中持续复用\n"
             "● 最优配置：lr=3×10⁻⁵, batch=32, γ=0.99, τ=0.005, ε_decay=0.998, PER α=0.6 β=0.4",
             Inches(0.5), Inches(3.9), Inches(12.1), Inches(2.0), title_color=ACCENT)

    # 训练效率
    add_card(slide, "训练增强机制效果",
             "基础DDQN→+PER→+SJF预填充→+SJF偏置探索→+GA-HPO：收敛episode从280降至100\n"
             "SJF预填充加速收敛46%（280→150 episodes），SJF偏置探索进一步至120 episodes\n"
             "正式训练仅需~52分钟（30个epoch），满足实际工程应用需求",
             Inches(0.5), Inches(5.6), Inches(12.1), Inches(1.3), title_color=ACCENT2)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 14: 消融实验 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "四、消融实验", Inches(0.5), Inches(0.25), Inches(8), Inches(0.6),
                    font_size=28, color=PRIMARY)

    # 消融实验柱状图
    chart_data4 = CategoryChartData()
    chart_data4.categories = ['完整模型', '-动作掩码', '-GA-HPO', '-交叉注意力', '-Dueling', '-PER', '-SJF预填充', '-GraphTrans', '-SJF偏置']
    chart_data4.add_series('Makespan变化(%)', (0, 13.0, 12.7, 10.5, 8.0, 5.0, 4.0, 3.0, 2.1))

    chart4 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.5), Inches(1.1), Inches(7.5), Inches(3.5),
        chart_data4
    ).chart
    chart4.has_legend = False
    chart4.chart_title.has_text_frame = True
    chart4.chart_title.text_frame.paragraphs[0].text = "去除各组件后Makespan变化 (%↑越大=越重要)"
    chart4.chart_title.text_frame.paragraphs[0].font.size = Pt(12)
    chart4.chart_title.text_frame.paragraphs[0].font.color.rgb = DARK
    plot4 = chart4.plots[0]
    series4 = plot4.series[0]
    series4.format.fill.solid()
    series4.format.fill.fore_color.rgb = ACCENT3

    # 分析
    add_card(slide, "消融实验结论",
             "● 影响最大三个组件：动作掩码(+13%)、GA-HPO(+12.7%)、交叉注意力(+10.5%)\n"
             "● 动作掩码不仅是约束保证，更是性能关键支撑\n"
             "● 交叉注意力是双流架构的核心——实现任务需求与资源能力的精准匹配\n"
             "● PER(+5%)和SJF预填充(+4%)主要贡献在训练效率和收敛速度\n"
             "● Graph Transformer(+3%)在复杂DAG工作流上贡献更显著(约+5.2%)\n"
             "● 所有组件对整体性能均有正向贡献，验证了算法设计的合理性",
             Inches(8.3), Inches(1.1), Inches(4.5), Inches(5.5), title_color=ACCENT)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 15: 案例研究 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "四、实际工作流案例研究", Inches(0.5), Inches(0.25), Inches(10), Inches(0.6),
                    font_size=28, color=PRIMARY)

    # 案例1：监管报送
    add_card(slide, "案例1：保单登记股份增量正式报送调度",
             "● 业务场景：向金融监管机构定期报送保单数据\n"
             "● 工作流规模：49个任务，DAG深度15层，最大宽度8\n"
             "● 原始Makespan：14,238秒（约3小时57分钟）\n"
             "● FE-IDDQN优化：13,845秒 → 改进2.76%（节省393秒）\n"
             "● 传统启发式仅改进0.04%-0.06%（依赖约束限制并行空间）\n"
             "● 34个省级子流程累积效应显著，依赖满足率100%",
             Inches(0.5), Inches(1.1), Inches(5.8), Inches(3.2), title_color=ACCENT)

    # 案例2：反洗钱检核
    add_card(slide, "案例2：EAST2-检核（反洗钱合规检核）",
             "● 业务场景：保全、理赔等业务数据的合规性批处理检核\n"
             "● 工作流规模：10个任务，DAG深度5层，单任务长耗时\n"
             "● 原始Makespan：42,232秒（约11小时43分钟）\n"
             "● FE-IDDQN优化：35,780秒 → 改进15.28%（节省6,452秒）\n"
             "● FIFO反而恶化-3.11%（串行分配到单一节点）\n"
             "● SJF/EFT几乎无改进（0.01%），FE-IDDQN学会异构匹配",
             Inches(6.8), Inches(1.1), Inches(5.8), Inches(3.2), title_color=ACCENT2)

    # 总结
    add_card(slide, "案例研究总结",
             "● 监管报送（大规模复杂DAG）：FE-IDDQN通过Graph Transformer感知DAG结构，优先优化关键路径\n"
             "● 反洗钱检核（少任务长耗时）：FE-IDDQN通过学习任务特征实现精准的任务-资源匹配\n"
             "● 两个案例场景迥异，FE-IDDQN均取得最优结果 → 验证算法的通用性和实用性\n"
             "● 100%依赖约束满足率 → 在金融监管合规场景中尤为关键",
             Inches(0.5), Inches(4.6), Inches(12.1), Inches(2.2), title_color=GOLD)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 16: 总结 - 主要贡献 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "五、总结 — 主要贡献", Inches(0.5), Inches(0.25), Inches(8), Inches(0.6),
                    font_size=28, color=PRIMARY)

    contributions = [
        ("贡献一：特征增强双流网络架构",
         "任务流+资源流独立编码 → 交叉注意力融合 → Dueling输出\n"
         "基于DolphinScheduler领域知识的16+7维特征体系\n"
         "交叉注意力贡献+10.5%性能提升"),

        ("贡献二：SJF预填充+偏置探索训练加速",
         "SJF启发式生成高质量种子经验 → 解决DRL冷启动问题\n"
         "SJF偏置ε-贪心 → 提高探索效率\n"
         "收敛速度提高46%（280→150 episodes）"),

        ("贡献三：GA两阶段超参数自动优化",
         "GA架构搜索 + Optuna训练参数搜索 → 消除人工经验依赖\n"
         "两阶段协同效应 → Makespan额外降低12.7%\n"
         "重复可用的一次性离线优化"),

        ("贡献四：基于生产数据的端到端实验框架",
         "首次使用金融行业DolphinScheduler生产数据（非合成数据）\n"
         "历史重放仿真器 → 端到端训练与评估\n"
         "监管报送和反洗钱检核真实案例验证"),
    ]

    for i, (title, desc) in enumerate(contributions):
        row = i // 2
        col = i % 2
        x = Inches(0.5 + col * 6.2)
        y = Inches(1.1 + row * 2.9)
        colors = [ACCENT, ACCENT2, GOLD, RGBColor(0x8E, 0x44, 0xAD)]
        add_card(slide, title, desc, x, y, Inches(5.8), Inches(2.5), title_color=colors[i])

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 17: 核心数据汇总 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "五、核心实验数据汇总", Inches(0.5), Inches(0.25), Inches(8), Inches(0.6),
                    font_size=28, color=PRIMARY)

    # 大数字展示
    metrics = [
        ("43.7%", "Makespan降低\n(vs FIFO)", ACCENT),
        ("9.4%", "Makespan降低\n(vs 最优启发式EFT)", ACCENT2),
        ("5.7%", "Makespan降低\n(vs 同族DDQN)", RGBColor(0x8E, 0x44, 0xAD)),
        ("58.6%", "资源利用率\n(最高)", GOLD),
        ("0.87", "负载均衡指数\n(最高)", ACCENT3),
        ("~5ms", "推理延迟\n(满足在线要求)", PRIMARY),
    ]

    for i, (number, label, color) in enumerate(metrics):
        col = i % 3
        row = i // 3
        x = Inches(0.8 + col * 4.2)
        y = Inches(1.2 + row * 2.8)

        # 大数字
        add_title_shape(slide, number, x, y, Inches(3.5), Inches(1.0),
                        font_size=48, color=color, bold=True, alignment=PP_ALIGN.CENTER)
        # 标签
        add_body_text(slide, label, x, y + Inches(1.1), Inches(3.5), Inches(0.8),
                      font_size=14, color=DARK, alignment=PP_ALIGN.CENTER, line_spacing=1.3)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 18: 展望 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_top_bar(slide, prs)

    add_title_shape(slide, "五、不足与展望", Inches(0.5), Inches(0.25), Inches(8), Inches(0.6),
                    font_size=28, color=PRIMARY)

    future_items = [
        ("数据多样性与跨领域泛化", "引入公开基准数据集和其他行业数据，验证算法泛化能力"),
        ("动态环境在线学习", "研究在线增量学习策略，应对任务失败、资源故障等动态事件"),
        ("大规模资源池扩展", "层次化动作空间设计、多智能体RL，支持数十至数百节点集群"),
        ("多目标Pareto学习", "从加权聚合到Pareto最优策略集，引入能耗、成本、QoS等目标"),
        ("超参数优化效率", "引入ENAS/DARTS等权重共享NAS和元学习方法，提升1-2个数量级"),
        ("与DolphinScheduler深度集成", "作为调度插件替代内置策略，解决模型在线更新和A/B测试"),
        ("可解释性增强", "注意力权重可视化、策略蒸馏与规则提取，提高透明度和可信度"),
    ]

    for i, (title, desc) in enumerate(future_items):
        row = i // 2
        col = i % 2
        if i < 6:
            x = Inches(0.5 + col * 6.2)
            y = Inches(1.1 + row * 1.55)
            add_card(slide, title, desc, x, y, Inches(5.8), Inches(1.3),
                     title_color=ACCENT if col == 0 else ACCENT2)
        else:
            # 最后一项居中
            add_card(slide, title, desc, Inches(0.5), Inches(1.1 + 3 * 1.55), Inches(12.1), Inches(1.3),
                     title_color=GOLD)

    add_bottom_bar(slide, prs)
    add_slide_number(slide, prs, slide_num)

    # ==================== Slide 19: 致谢 ====================
    slide_num += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, PRIMARY)

    add_title_shape(slide, "谢谢各位老师！", Inches(0.5), Inches(1.5), Inches(12), Inches(1.5),
                    font_size=48, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    deco_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.2), Inches(2.3), Inches(0.04))
    deco_line.fill.solid()
    deco_line.fill.fore_color.rgb = ACCENT
    deco_line.line.fill.background()

    add_body_text(slide, "恳请各位老师批评指正", Inches(0.5), Inches(3.6), Inches(12), Inches(0.8),
                  font_size=24, color=RGBColor(0xBB, 0xCC, 0xDD), alignment=PP_ALIGN.CENTER)

    add_body_text(slide, "答辩人：XXX\n2026年X月", Inches(0.5), Inches(4.5), Inches(12), Inches(1.2),
                  font_size=18, color=RGBColor(0xAA, 0xBB, 0xCC), alignment=PP_ALIGN.CENTER, line_spacing=1.6)

    # ==================== 保存 ====================
    output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "1", "答辩PPT_GA-HPO_FE-IDDQN.pptx")
    prs.save(output_path)
    print(f"✅ 答辩PPT已生成：{output_path}")
    print(f"📊 共 {slide_num} 页幻灯片")


if __name__ == "__main__":
    create_presentation()
